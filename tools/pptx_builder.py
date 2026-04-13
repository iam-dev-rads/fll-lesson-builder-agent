import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configure logging
logger = logging.getLogger(__name__)

# FLL Branded Colors
FLL_BLUE = RGBColor(0, 102, 204)
FLL_ORANGE = RGBColor(255, 102, 0)

def _apply_title_style(shape):
    """
    Applies the FLL brand style to a slide title.
    """
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = FLL_BLUE

def _apply_body_style(shape, size=Pt(24)):
    """
    Applies consistent font sizing to body text.
    """
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        paragraph.font.size = size

def build_presentation(content: dict, output_folder: str) -> str:
    """
    Builds a PowerPoint presentation based on the FLL lesson content JSON.
    Receives validated content dict and output folder path.
    Returns the absolute path to the saved .pptx file.
    """
    prs = Presentation()
    
    # Set slide dimensions to 13.33" x 7.5" (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide Layout Indices
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_AND_CONTENT = 1
    
    # 1. Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = content["lesson_title"]
    subtitle.text = "FLL Team Lesson"
    _apply_title_style(title)
    
    # 2. Slide 2: Opening Hook
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Let's Start With A Question!"
    body.text = content["opening_hook"]
    _apply_title_style(title)
    _apply_body_style(body)
    
    # 3. Slides 3 to N: Content Slides
    recap_bullets = []
    for slide_data in content["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT])
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = slide_data["title"]
        _apply_title_style(title)
        
        # Add main bullet points
        tf = body.text_frame
        tf.text = "" # Clear the default placeholder text
        for i, bullet in enumerate(slide_data["bullet_points"]):
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
            # Collect the first bullet of each slide for the recap
            if i == 0:
                recap_bullets.append(bullet)
        
        # Add Example text box at the bottom
        left = Inches(1)
        top = Inches(6.2)
        width = Inches(11.3)
        height = Inches(1)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        p = tf.add_paragraph()
        
        # "Example:" Label in FLL Orange
        run = p.add_run()
        run.text = "Example: "
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = FLL_ORANGE
        
        # The actual example text
        run = p.add_run()
        run.text = slide_data["example"]
        run.font.size = Pt(18)
        run.font.italic = True

    # 4. Slide N+1: Let's Try It!
    exercise = content["closing_exercise"]
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = f"Let's Try It: {exercise['title']}"
    body.text = exercise["instructions"]
    _apply_title_style(title)
    _apply_body_style(body)
    
    # Small box for expected outcome
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(1))
    tf = tx_box.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"What success looks like: {exercise['expected_outcome']}"
    run.font.size = Pt(18)
    run.font.italic = True
    run.font.color.rgb = FLL_BLUE
    
    # 5. Slide N+2: Recap Slide
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "What Did We Learn Today?"
    _apply_title_style(title)
    
    tf = body.text_frame
    tf.text = ""
    for bullet in recap_bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)

    # Save the presentation
    output_path = os.path.join(output_folder, "lesson_slides.pptx")
    abs_path = os.path.abspath(output_path)
    prs.save(abs_path)
    
    logger.info(f"PowerPoint saved to: {abs_path}")
    return abs_path
